from pptx.shapes.base import BaseShape
def is_valid_shape(shape: BaseShape) -> bool:
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        if shape.placeholder_format and shape.placeholder_format.type:
            placeholder_type = (
                str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]  
            )
            if placeholder_type == "SLIDE_NUMBER":
                return False
            if placeholder_type == "FOOTER" and text.isdigit():
                return False

    return True